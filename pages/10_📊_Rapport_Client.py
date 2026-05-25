"""
📊 RAPPORT CLIENT AUTOMATISÉ
==============================
Génère un bilan SEO mensuel professionnel (.pptx) au format Uplix.
Combine Ahrefs API v3 + GSC API + Claude pour produire un deck
prêt à présenter : KPIs, graphiques, top pages, recommandations.
Importable dans Google Slides.
"""

import streamlit as st
import requests
import base64
import pandas as pd
import plotly.graph_objects as go
import anthropic
import json
import re
import io
import os
import logging
from datetime import datetime, timedelta
from urllib.parse import urlparse
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker

from utils.auth import check_password

logger = logging.getLogger(__name__)

# =============================================================================
# PAGE CONFIG
# =============================================================================
st.set_page_config(
    page_title="Rapport Client | Ma Toolbox SEO",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

check_password()

st.markdown("""
<style>
    .stMetric > div { padding: 8px; }
    .section-title {
        font-size: 1.5rem; font-weight: 600; color: #1E3A5F;
        border-left: 4px solid #4285F4; padding-left: 1rem;
        margin: 2rem 0 1rem 0;
    }
</style>
""", unsafe_allow_html=True)


# =============================================================================
# UPLIX DA CONSTANTS
# =============================================================================
UPLIX_BLACK = RGBColor(0x00, 0x00, 0x00)
UPLIX_DARK_GRAY = RGBColor(0x59, 0x59, 0x59)
UPLIX_LIGHT_GRAY = RGBColor(0xEE, 0xEE, 0xEE)
UPLIX_BG = RGBColor(0xF8, 0xF8, 0xF8)
UPLIX_BLUE = RGBColor(0x42, 0x85, 0xF4)
UPLIX_ORANGE = RGBColor(0xFF, 0xAB, 0x40)
UPLIX_TEAL = RGBColor(0x00, 0x97, 0xA7)
UPLIX_RED = RGBColor(0xEF, 0x44, 0x44)
UPLIX_GREEN = RGBColor(0x22, 0xC5, 0x5E)
UPLIX_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_TITLE = "Poppins"
FONT_BODY = "Poppins Light"

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Chart color palette (matplotlib hex)
CHART_BLUE = "#4285F4"
CHART_ORANGE = "#FFAB40"
CHART_TEAL = "#0097A7"
CHART_RED = "#EF4444"
CHART_GRAY = "#78909C"


# =============================================================================
# CLIENTS
# =============================================================================
CLIENTS_DIR = os.path.join(os.path.dirname(__file__), "clients")


def load_clients() -> dict:
    clients = {}
    if os.path.isdir(CLIENTS_DIR):
        for fname in sorted(os.listdir(CLIENTS_DIR)):
            if fname.endswith(".json"):
                try:
                    with open(os.path.join(CLIENTS_DIR, fname), "r", encoding="utf-8") as f:
                        data = json.load(f)
                    name = data.get("name", fname.replace(".json", ""))
                    if data.get("domain"):
                        clients[name] = data
                except Exception:
                    pass
    return clients


# =============================================================================
# API CLIENTS
# =============================================================================
AHREFS_API_BASE = "https://api.ahrefs.com/v3"


class AhrefsAPI:
    def __init__(self, token: str):
        self.session = requests.Session()
        self.session.headers.update({
            "Authorization": f"Bearer {token}",
            "Accept": "application/json",
        })

    def _get(self, endpoint: str, params: dict):
        params["output"] = "json"
        try:
            r = self.session.get(f"{AHREFS_API_BASE}/{endpoint}", params=params, timeout=60)
            r.raise_for_status()
            return r.json()
        except Exception as e:
            logger.error(f"Ahrefs [{endpoint}]: {e}")
            return None

    @staticmethod
    def _rows(data):
        if data is None:
            return []
        if isinstance(data, list):
            return data
        if isinstance(data, dict):
            for v in data.values():
                if isinstance(v, list):
                    return v
        return []

    def domain_rating(self, target, date):
        return self._get("site-explorer/domain-rating", {"target": target, "date": date})

    def metrics(self, target, date, mode="subdomains"):
        return self._get("site-explorer/metrics", {"target": target, "date": date, "mode": mode})

    def backlinks_stats(self, target, date, mode="subdomains"):
        return self._get("site-explorer/backlinks-stats", {"target": target, "date": date, "mode": mode})

    def metrics_history(self, target, date_from, date_to, mode="subdomains"):
        return self._get("site-explorer/metrics-history", {
            "target": target, "date_from": date_from, "date_to": date_to,
            "mode": mode, "history_grouping": "monthly",
            "select": "date,org_traffic,org_cost,org_keywords,paid_traffic",
        })

    def top_pages(self, target, date, country=None, limit=15, mode="subdomains"):
        p = {
            "target": target, "date": date, "mode": mode,
            "select": "url,sum_traffic,top_keyword,top_keyword_best_position,keywords,value",
            "order_by": "sum_traffic:desc", "limit": limit,
        }
        if country:
            p["country"] = country
        return self._get("site-explorer/top-pages", p)

    def organic_keywords(self, target, date, country="FR", limit=30, mode="subdomains"):
        return self._get("site-explorer/organic-keywords", {
            "target": target, "date": date, "mode": mode,
            "country": country, "limit": limit,
            "select": "keyword,volume,best_position,best_position_url,sum_traffic,keyword_difficulty,best_position_diff",
            "order_by": "sum_traffic:desc",
        })


class GSCAPI:
    def __init__(self):
        self.service = None
        try:
            from google.oauth2.credentials import Credentials
            from googleapiclient.discovery import build
            cid = st.secrets.get("GSC_CLIENT_ID", "")
            cs = st.secrets.get("GSC_CLIENT_SECRET", "")
            rt = st.secrets.get("GSC_REFRESH_TOKEN", "")
            if cid and cs and rt:
                creds = Credentials(token=None, refresh_token=rt, client_id=cid,
                                    client_secret=cs, token_uri="https://oauth2.googleapis.com/token")
                self.service = build("searchconsole", "v1", credentials=creds)
                return
            if "GSC_SERVICE_ACCOUNT" in st.secrets:
                from google.oauth2 import service_account
                creds = service_account.Credentials.from_service_account_info(
                    dict(st.secrets["GSC_SERVICE_ACCOUNT"]),
                    scopes=["https://www.googleapis.com/auth/webmasters.readonly"])
                self.service = build("searchconsole", "v1", credentials=creds)
        except Exception as e:
            logger.warning(f"GSC init: {e}")

    @property
    def ok(self):
        return self.service is not None

    def totals(self, site, start, end):
        if not self.service:
            return None
        try:
            r = self.service.searchanalytics().query(
                siteUrl=site, body={"startDate": start, "endDate": end, "dataState": "final"}
            ).execute()
            rows = r.get("rows", [])
            return rows[0] if rows else None
        except Exception:
            return None

    def queries(self, site, start, end, limit=20):
        if not self.service:
            return []
        try:
            r = self.service.searchanalytics().query(
                siteUrl=site,
                body={"startDate": start, "endDate": end, "dimensions": ["query"],
                      "rowLimit": limit, "dataState": "final"}
            ).execute()
            return r.get("rows", [])
        except Exception:
            return []


# =============================================================================
# UTILS
# =============================================================================
def clean_domain(url):
    url = url.strip().lower()
    if not url.startswith(("http://", "https://")):
        url = "https://" + url
    parsed = urlparse(url)
    d = parsed.netloc or parsed.path
    return re.sub(r'^www\.', '', d).rstrip('/')


def fmt(n):
    if n is None:
        return "N/A"
    n = float(n)
    if n >= 1_000_000:
        return f"{n/1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n/1_000:.1f}K"
    return f"{n:.0f}"


def delta_pct(cur, prev):
    if not prev:
        return "N/A", "neutral"
    diff = cur - prev
    pct = (diff / abs(prev)) * 100
    sign = "+" if diff >= 0 else ""
    direction = "up" if diff > 0 else "down" if diff < 0 else "neutral"
    return f"{sign}{pct:.1f}%", direction


def extract_metric(data, key, default=0):
    if data is None:
        return default
    if isinstance(data, dict):
        if key in data:
            return data[key] or default
        for v in data.values():
            if isinstance(v, dict) and key in v:
                return v[key] or default
            if isinstance(v, list) and v and isinstance(v[0], dict) and key in v[0]:
                return v[0][key] or default
    return default


# =============================================================================
# CHART GENERATION (Matplotlib → PNG buffer)
# =============================================================================
def chart_traffic_evolution(history: list) -> io.BytesIO:
    plt.rcParams.update({"font.family": "sans-serif", "font.size": 10})
    fig, ax = plt.subplots(figsize=(8, 3.2), dpi=150)

    dates = [h["date"] for h in history]
    traffic = [h.get("traffic", 0) for h in history]

    ax.fill_between(range(len(dates)), traffic, alpha=0.15, color=CHART_BLUE)
    ax.plot(range(len(dates)), traffic, color=CHART_BLUE, linewidth=2.5, marker="o", markersize=5)

    ax.set_xticks(range(len(dates)))
    ax.set_xticklabels(dates, rotation=45, ha="right", fontsize=8)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: fmt(x)))
    ax.set_ylabel("Trafic organique estimé", fontsize=9, color="#595959")
    ax.grid(axis="y", alpha=0.3)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_keywords_evolution(history: list) -> io.BytesIO:
    plt.rcParams.update({"font.family": "sans-serif", "font.size": 10})
    fig, ax = plt.subplots(figsize=(8, 3.2), dpi=150)

    dates = [h["date"] for h in history]
    keywords = [h.get("keywords", 0) for h in history]

    ax.fill_between(range(len(dates)), keywords, alpha=0.15, color=CHART_TEAL)
    ax.plot(range(len(dates)), keywords, color=CHART_TEAL, linewidth=2.5, marker="o", markersize=5)

    ax.set_xticks(range(len(dates)))
    ax.set_xticklabels(dates, rotation=45, ha="right", fontsize=8)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: fmt(x)))
    ax.set_ylabel("Mots-clés organiques", fontsize=9, color="#595959")
    ax.grid(axis="y", alpha=0.3)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


# =============================================================================
# PPTX GENERATION — UPLIX DA
# =============================================================================

def _set_font(run, name=FONT_BODY, size=Pt(10), color=UPLIX_BLACK, bold=False):
    run.font.name = name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold


def _add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                 font_size=Pt(10), color=UPLIX_BLACK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, font_name, font_size, color, bold)
    return txBox


def _add_uplix_footer(slide):
    _add_textbox(slide, Inches(8.8), Inches(5.15), Inches(1), Inches(0.35),
                 "Uplix", FONT_TITLE, Pt(10), UPLIX_DARK_GRAY, bold=True,
                 alignment=PP_ALIGN.RIGHT)


def _add_source_footer(slide, source="Ahrefs"):
    _add_textbox(slide, Inches(3.5), Inches(5.15), Inches(3), Inches(0.35),
                 f"Source : {source}", FONT_BODY, Pt(8), UPLIX_DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)


def _add_slide_title(slide, title_text):
    _add_textbox(slide, Inches(0.4), Inches(0.2), Inches(8), Inches(0.5),
                 title_text, FONT_TITLE, Pt(18), UPLIX_BLACK, bold=True)


def _add_content_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = UPLIX_LIGHT_GRAY
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_kpi_block(slide, left, top, width, label, value, subtitle="", color=UPLIX_BLACK):
    _add_textbox(slide, left, top, width, Inches(0.25),
                 label.upper(), FONT_BODY, Pt(7), UPLIX_DARK_GRAY, bold=True)
    _add_textbox(slide, left, top + Inches(0.25), width, Inches(0.45),
                 str(value), FONT_TITLE, Pt(24), color, bold=True)
    if subtitle:
        _add_textbox(slide, left, top + Inches(0.7), width, Inches(0.3),
                     subtitle, FONT_BODY, Pt(8), UPLIX_DARK_GRAY)


def create_cover_slide(prs, client_name, period, logo_bytes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = UPLIX_WHITE

    # Left vertical accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(1.8), Inches(0.06), Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = UPLIX_BLACK
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.6), Inches(1.8), Inches(5), Inches(0.6),
                 "BILAN SEO MENSUEL", FONT_TITLE, Pt(26), UPLIX_BLACK, bold=True)
    _add_textbox(slide, Inches(0.6), Inches(2.4), Inches(5), Inches(0.5),
                 client_name.upper(), FONT_TITLE, Pt(18), UPLIX_DARK_GRAY, bold=False)
    _add_textbox(slide, Inches(0.6), Inches(2.9), Inches(5), Inches(0.4),
                 period, FONT_BODY, Pt(11), UPLIX_DARK_GRAY)

    if logo_bytes:
        logo_stream = io.BytesIO(logo_bytes)
        try:
            img = Image.open(logo_stream)
            w, h = img.size
            max_h = Inches(1.2)
            max_w = Inches(2.5)
            ratio = min(max_w / Emu(int(w * 914400 / 96)), max_h / Emu(int(h * 914400 / 96)))
            final_w = int(w * 914400 / 96 * ratio)
            final_h = int(h * 914400 / 96 * ratio)
            logo_stream.seek(0)
            slide.shapes.add_picture(logo_stream, Inches(6.5), Inches(1.8), final_w, final_h)
        except Exception:
            pass

    _add_textbox(slide, Inches(0.4), Inches(5.0), Inches(3), Inches(0.35),
                 "Uplix · Tous droits réservés", FONT_BODY, Pt(8), UPLIX_DARK_GRAY)


def create_section_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = UPLIX_WHITE

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(2.3), Inches(0.06), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = UPLIX_BLACK
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.6), Inches(2.3), Inches(7), Inches(0.7),
                 title_text.upper(), FONT_TITLE, Pt(24), UPLIX_BLACK, bold=True)

    _add_textbox(slide, Inches(0.4), Inches(5.0), Inches(2), Inches(0.35),
                 "Uplix", FONT_TITLE, Pt(10), UPLIX_DARK_GRAY, bold=True)


def create_kpi_slide(prs, kpis: list):
    """kpis = list of (label, value, subtitle, color)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = UPLIX_WHITE

    _add_slide_title(slide, "KPIs du mois")

    _add_content_card(slide, Inches(0.4), Inches(0.9), Inches(9.2), Inches(3.8))

    n = len(kpis)
    card_w = 9.0 / n
    for i, (label, value, subtitle, color) in enumerate(kpis):
        left = Inches(0.5 + i * card_w)
        _add_kpi_block(slide, left, Inches(1.3), Inches(card_w - 0.1),
                       label, value, subtitle, color)

        if i < n - 1:
            sep = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.45 + (i + 1) * card_w), Inches(1.2),
                Inches(0.01), Inches(1.5)
            )
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            sep.line.fill.background()

    _add_uplix_footer(slide)
    _add_source_footer(slide, "Ahrefs + GSC")


def create_chart_slide(prs, title, chart_buf, source="Ahrefs"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = UPLIX_WHITE

    _add_slide_title(slide, title)
    _add_content_card(slide, Inches(0.4), Inches(0.85), Inches(9.2), Inches(4.0))

    slide.shapes.add_picture(chart_buf, Inches(0.6), Inches(0.95), Inches(8.8), Inches(3.6))

    _add_uplix_footer(slide)
    _add_source_footer(slide, source)


def create_table_slide(prs, title, headers, rows, source="Ahrefs", col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = UPLIX_WHITE

    _add_slide_title(slide, title)

    n_rows = min(len(rows), 12)
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows + 1, n_cols,
        Inches(0.4), Inches(0.85),
        Inches(9.2), Inches(0.35 * (n_rows + 1))
    )
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        _set_font(run, FONT_TITLE, Pt(8), UPLIX_WHITE, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x32, 0x33, 0x35)

    for ri, row in enumerate(rows[:n_rows]):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            _set_font(run, FONT_BODY, Pt(8), UPLIX_BLACK)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)

    _add_uplix_footer(slide)
    _add_source_footer(slide, source)


def create_backlinks_slide(prs, bl_data):
    kpis = [
        ("Domain Rating", str(bl_data.get("domain_rating", "N/A")), "", UPLIX_BLUE),
        ("Backlinks", fmt(bl_data.get("live_backlinks", 0)), "Liens entrants actifs", UPLIX_BLACK),
        ("Domaines référents", fmt(bl_data.get("live_refdomains", 0)), "Domaines uniques", UPLIX_TEAL),
        ("DoFollow", fmt(bl_data.get("dofollow", 0)), f"NoFollow : {fmt(bl_data.get('nofollow', 0))}", UPLIX_GREEN),
    ]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = UPLIX_WHITE

    _add_slide_title(slide, "Profil de liens")
    _add_content_card(slide, Inches(0.4), Inches(0.9), Inches(9.2), Inches(2.5))

    n = len(kpis)
    card_w = 9.0 / n
    for i, (label, value, subtitle, color) in enumerate(kpis):
        _add_kpi_block(slide, Inches(0.5 + i * card_w), Inches(1.2),
                       Inches(card_w - 0.1), label, value, subtitle, color)

    _add_uplix_footer(slide)
    _add_source_footer(slide, "Ahrefs")


def create_recommendations_slide(prs, reco_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = UPLIX_WHITE

    _add_slide_title(slide, "Synthèse & Recommandations")

    _add_content_card(slide, Inches(0.4), Inches(0.85), Inches(9.2), Inches(4.0))

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(8.8), Inches(3.7))
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = reco_text.strip().split("\n")
    first = True
    for line in lines:
        line = line.strip()
        if not line:
            continue

        if not first:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
            first = False

        p.space_after = Pt(4)

        if line.startswith("## ") or line.startswith("### "):
            clean = line.lstrip("#").strip()
            run = p.add_run()
            run.text = clean
            _set_font(run, FONT_TITLE, Pt(11), UPLIX_BLACK, bold=True)
            p.space_before = Pt(8)
        elif line.startswith("- ") or line.startswith("* "):
            run = p.add_run()
            run.text = "  •  " + line[2:]
            _set_font(run, FONT_BODY, Pt(9), UPLIX_DARK_GRAY)
        elif line.startswith("**") and line.endswith("**"):
            run = p.add_run()
            run.text = line.replace("**", "")
            _set_font(run, FONT_TITLE, Pt(10), UPLIX_BLACK, bold=True)
        else:
            run = p.add_run()
            run.text = line
            _set_font(run, FONT_BODY, Pt(9), UPLIX_DARK_GRAY)

    _add_uplix_footer(slide)


def create_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = UPLIX_WHITE

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(2.3), Inches(0.06), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = UPLIX_BLACK
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.6), Inches(2.3), Inches(7), Inches(0.7),
                 "MERCI", FONT_TITLE, Pt(28), UPLIX_BLACK, bold=True)
    _add_textbox(slide, Inches(0.6), Inches(3.0), Inches(7), Inches(0.4),
                 "Des questions ?", FONT_BODY, Pt(14), UPLIX_DARK_GRAY)

    _add_textbox(slide, Inches(0.4), Inches(5.0), Inches(3), Inches(0.35),
                 "Uplix · Tous droits réservés", FONT_BODY, Pt(8), UPLIX_DARK_GRAY)


def generate_report_pptx(report: dict) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover
    create_cover_slide(prs, report["client_name"], report["period"],
                       report.get("logo_bytes"))

    # 2. KPIs
    kpis = report.get("kpis_list", [])
    if kpis:
        create_kpi_slide(prs, kpis)

    # 3. Traffic evolution chart
    history = report.get("traffic_history", [])
    if history:
        buf = chart_traffic_evolution(history)
        create_chart_slide(prs, "Évolution du trafic organique", buf, "Ahrefs")

    # 4. Keywords evolution chart
    if history and any(h.get("keywords", 0) > 0 for h in history):
        buf = chart_keywords_evolution(history)
        create_chart_slide(prs, "Évolution des mots-clés organiques", buf, "Ahrefs")

    # 5. Top pages
    top_pages = report.get("top_pages", [])
    if top_pages:
        headers = ["URL", "Trafic", "Top mot-clé", "Pos.", "KWs"]
        rows = []
        for p in top_pages[:12]:
            url = p.get("url", "")
            if len(url) > 55:
                url = url[:52] + "..."
            rows.append([
                url, fmt(p.get("traffic", 0)),
                p.get("top_keyword", ""), str(p.get("position", "")),
                str(p.get("keywords_count", ""))
            ])
        create_table_slide(prs, "Top pages par trafic", headers, rows, "Ahrefs",
                          [4.5, 1.0, 2.0, 0.7, 0.7])

    # 6. Keywords
    keywords = report.get("keywords", [])
    if keywords:
        headers = ["Mot-clé", "Volume", "Pos.", "Mvt", "Trafic"]
        rows = []
        for kw in keywords[:12]:
            diff = kw.get("position_diff", 0) or 0
            if diff > 0:
                mvt = f"↑ +{diff}"
            elif diff < 0:
                mvt = f"↓ {diff}"
            else:
                mvt = "="
            rows.append([
                kw.get("keyword", ""), fmt(kw.get("volume", 0)),
                str(kw.get("position", "")), mvt, fmt(kw.get("traffic", 0))
            ])
        create_table_slide(prs, "Mots-clés principaux", headers, rows, "Ahrefs",
                          [3.5, 1.2, 0.8, 1.0, 1.2])

    # 7. Backlinks
    bl = report.get("backlinks", {})
    if bl:
        create_backlinks_slide(prs, bl)

    # 8. GSC
    gsc = report.get("gsc", {})
    if gsc.get("available"):
        cur = gsc.get("current", {})
        prev = gsc.get("previous", {})
        clicks_d, _ = delta_pct(cur.get("clicks", 0), prev.get("clicks", 0))
        impr_d, _ = delta_pct(cur.get("impressions", 0), prev.get("impressions", 0))

        gsc_kpis = [
            ("Clics", fmt(cur.get("clicks", 0)), clicks_d, UPLIX_BLUE),
            ("Impressions", fmt(cur.get("impressions", 0)), impr_d, UPLIX_BLACK),
            ("CTR moyen", f"{cur.get('ctr', 0) * 100:.2f}%", "", UPLIX_TEAL),
            ("Position moyenne", f"{cur.get('position', 0):.1f}", "", UPLIX_ORANGE),
        ]
        create_kpi_slide(prs, gsc_kpis)

        # GSC top queries table
        queries = gsc.get("top_queries", [])
        if queries:
            headers = ["Requête", "Clics", "Impressions", "CTR", "Pos."]
            rows = []
            for q in queries[:12]:
                keys = q.get("keys", [""])
                rows.append([
                    keys[0] if keys else "",
                    fmt(q.get("clicks", 0)), fmt(q.get("impressions", 0)),
                    f"{q.get('ctr', 0) * 100:.1f}%", f"{q.get('position', 0):.1f}"
                ])
            create_table_slide(prs, "Top requêtes GSC", headers, rows, "Google Search Console",
                              [3.5, 1.2, 1.5, 1.0, 1.0])

    # 9. Recommendations
    reco = report.get("recommendations", "")
    if reco:
        create_recommendations_slide(prs, reco)

    # 10. Closing
    create_closing_slide(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# =============================================================================
# CLAUDE RECOMMENDATIONS
# =============================================================================
def generate_recommendations(api_key, data):
    client = anthropic.Anthropic(api_key=api_key)
    prompt = f"""Tu es un consultant SEO senior chez Uplix. Rédige la synthèse du bilan mensuel pour le client "{data['client_name']}" ({data['domain']}).

Données du mois :
- Trafic organique : {data.get('current_traffic', 'N/A')} ({data.get('traffic_delta', 'N/A')})
- Mots-clés : {data.get('current_keywords', 'N/A')} ({data.get('keywords_delta', 'N/A')})
- DR : {data.get('domain_rating', 'N/A')} | Refdomains : {data.get('refdomains', 'N/A')}

Top pages :
{data.get('top_pages_summary', 'N/A')}

Mots-clés (mouvements) :
{data.get('keywords_summary', 'N/A')}

GSC : Clics {data.get('gsc_clicks', 'N/A')} | Impressions {data.get('gsc_impressions', 'N/A')} | CTR {data.get('gsc_ctr', 'N/A')} | Pos moy {data.get('gsc_position', 'N/A')}

Rédige en français, format markdown :
## Synthèse du mois
3-4 phrases factuelles

## Points positifs
- 2-3 bullets

## Points d'attention
- 2-3 bullets

## Actions recommandées
- 3-5 bullets priorisées et concrètes

Sois concis, factuel, orienté action. Pas de blabla.
"""
    r = client.messages.create(
        model="claude-sonnet-4-20250514", max_tokens=1500,
        messages=[{"role": "user", "content": prompt}]
    )
    return r.content[0].text


# =============================================================================
# LOCATIONS
# =============================================================================
LOCATIONS = {
    "France": {"code": 2250, "lang": "fr", "ahrefs": "FR"},
    "Belgique": {"code": 2056, "lang": "fr", "ahrefs": "BE"},
    "Suisse": {"code": 2756, "lang": "fr", "ahrefs": "CH"},
    "Canada (FR)": {"code": 2124, "lang": "fr", "ahrefs": "CA"},
    "États-Unis": {"code": 2840, "lang": "en", "ahrefs": "US"},
    "Royaume-Uni": {"code": 2826, "lang": "en", "ahrefs": "GB"},
    "Allemagne": {"code": 2276, "lang": "de", "ahrefs": "DE"},
    "Espagne": {"code": 2724, "lang": "es", "ahrefs": "ES"},
}


# =============================================================================
# STREAMLIT UI
# =============================================================================
st.title("📊 Rapport Client Automatisé")
st.markdown("*Génère un bilan SEO mensuel .pptx au format Uplix. Importable dans Google Slides.*")

with st.sidebar:
    st.header("⚙️ Configuration API")

    anthropic_key = st.text_input("Clé API Claude", value=st.secrets.get("ANTHROPIC_API_KEY", ""), type="password")
    ahrefs_token = st.text_input("Token Ahrefs", value=st.secrets.get("AHREFS_API_TOKEN", ""), type="password")

    st.divider()
    st.header("📊 Paramètres du rapport")

    clients = load_clients()
    client_names = list(clients.keys())
    if client_names:
        selected_client = st.selectbox("Client", ["— Saisie manuelle —"] + client_names)
    else:
        selected_client = "— Saisie manuelle —"

    if selected_client == "— Saisie manuelle —":
        target_domain = st.text_input("Domaine", placeholder="monsite.fr")
        client_display_name = st.text_input("Nom du client", placeholder="Mon Client")
    else:
        client_data = clients[selected_client]
        target_domain = client_data.get("domain", "")
        client_display_name = client_data.get("name", selected_client)
        st.markdown(f"**Domaine** : `{target_domain}`")

    selected_location = st.selectbox("Pays cible", list(LOCATIONS.keys()), index=0)

    report_period = st.selectbox("Période", ["Dernier mois", "3 derniers mois", "6 derniers mois"])

    logo_file = st.file_uploader("Logo du client (PNG/JPG)", type=["png", "jpg", "jpeg", "svg"])

    gsc_property = st.text_input("Propriété GSC (optionnel)", placeholder="sc-domain:monsite.fr")

    st.divider()
    st.markdown("""
    ### Le rapport contient
    1. **Cover** avec logo client
    2. **KPIs** avec évolutions
    3. **Courbe trafic** (12 mois)
    4. **Courbe mots-clés** (12 mois)
    5. **Top pages** (tableau)
    6. **Mots-clés principaux** (tableau)
    7. **Profil de liens**
    8. **GSC** (si connecté)
    9. **Recommandations** (Claude)
    10. **Slide de clôture**
    """)

# Validations
if not anthropic_key:
    st.warning("Configure ta clé API Claude dans la sidebar.")
    st.stop()
if not ahrefs_token:
    st.warning("Configure ton token Ahrefs dans la sidebar.")
    st.stop()
if not target_domain:
    st.info("Sélectionne un client ou entre un domaine.")
    st.stop()

domain = clean_domain(target_domain)
st.markdown(f"**Client** : {client_display_name} | **Domaine** : `{domain}` | **Pays** : {selected_location}")

if st.button("🚀 Générer le bilan", type="primary", use_container_width=True):
    loc = LOCATIONS[selected_location]
    ahrefs = AhrefsAPI(ahrefs_token)

    today = datetime.now()
    end_date = (today - timedelta(days=3)).strftime("%Y-%m-%d")
    period_days = {"Dernier mois": 30, "3 derniers mois": 90, "6 derniers mois": 180}[report_period]
    start_date = (today - timedelta(days=period_days + 3)).strftime("%Y-%m-%d")
    prev_end = start_date
    prev_start = (today - timedelta(days=period_days * 2 + 3)).strftime("%Y-%m-%d")
    history_start = (today - timedelta(days=365)).strftime("%Y-%m-%d")
    period_label = f"{(today - timedelta(days=period_days + 3)).strftime('%d/%m/%Y')} — {(today - timedelta(days=3)).strftime('%d/%m/%Y')}"

    logo_bytes = logo_file.read() if logo_file else None

    with st.status("Génération du bilan en cours...", expanded=True) as status:

        # ── Ahrefs ──
        st.write("📊 **Phase 1** — Collecte Ahrefs...")

        dr_data = ahrefs.domain_rating(domain, end_date)
        domain_rating = extract_metric(dr_data, "domain_rating")
        st.write(f"  ✅ DR : {domain_rating or 'N/A'}")

        cur_metrics = ahrefs.metrics(domain, end_date)
        current_traffic = extract_metric(cur_metrics, "org_traffic")
        current_keywords = extract_metric(cur_metrics, "org_keywords")

        prev_metrics = ahrefs.metrics(domain, prev_end)
        prev_traffic = extract_metric(prev_metrics, "org_traffic")
        prev_keywords = extract_metric(prev_metrics, "org_keywords")

        traffic_delta, t_dir = delta_pct(current_traffic, prev_traffic)
        kw_delta, k_dir = delta_pct(current_keywords, prev_keywords)

        st.write(f"  ✅ Trafic : {fmt(current_traffic)} ({traffic_delta})")
        st.write(f"  ✅ Mots-clés : {fmt(current_keywords)} ({kw_delta})")

        bl_raw = ahrefs.backlinks_stats(domain, end_date)
        bl_data = {
            "domain_rating": domain_rating,
            "live_backlinks": extract_metric(bl_raw, "live"),
            "live_refdomains": extract_metric(bl_raw, "live_refdomains"),
            "dofollow": extract_metric(bl_raw, "dofollow"),
            "nofollow": extract_metric(bl_raw, "nofollow"),
        }

        hist_raw = ahrefs.metrics_history(domain, history_start, end_date)
        traffic_history = []
        for row in AhrefsAPI._rows(hist_raw):
            traffic_history.append({
                "date": (row.get("date", "") or "")[:7],
                "traffic": row.get("org_traffic", 0) or 0,
                "keywords": row.get("org_keywords", 0) or 0,
                "cost": row.get("org_cost", 0) or 0,
            })
        st.write(f"  ✅ Historique : {len(traffic_history)} mois")

        tp_raw = ahrefs.top_pages(domain, end_date, country=loc["ahrefs"])
        top_pages = [{
            "url": r.get("url", ""), "traffic": r.get("sum_traffic", 0) or 0,
            "top_keyword": r.get("top_keyword", ""),
            "position": r.get("top_keyword_best_position", 0) or 0,
            "keywords_count": r.get("keywords", 0) or 0,
        } for r in AhrefsAPI._rows(tp_raw)]
        st.write(f"  ✅ Top pages : {len(top_pages)}")

        kw_raw = ahrefs.organic_keywords(domain, end_date, country=loc["ahrefs"])
        keywords_list = [{
            "keyword": r.get("keyword", ""), "volume": r.get("volume", 0) or 0,
            "position": r.get("best_position", 0) or 0,
            "position_diff": r.get("best_position_diff", 0) or 0,
            "traffic": r.get("sum_traffic", 0) or 0,
        } for r in AhrefsAPI._rows(kw_raw)]
        st.write(f"  ✅ Mots-clés : {len(keywords_list)}")

        # ── GSC ──
        st.write("📈 **Phase 2** — Google Search Console...")
        gsc = GSCAPI()
        gsc_report = {"available": False}
        if gsc.ok and gsc_property:
            cur_gsc = gsc.totals(gsc_property, start_date, end_date)
            prev_gsc = gsc.totals(gsc_property, prev_start, prev_end)
            queries_gsc = gsc.queries(gsc_property, start_date, end_date)
            if cur_gsc:
                gsc_report = {"available": True, "current": cur_gsc,
                              "previous": prev_gsc or {}, "top_queries": queries_gsc}
                st.write(f"  ✅ GSC — Clics : {fmt(cur_gsc.get('clicks', 0))}")
            else:
                st.write("  ⚠️ Pas de données GSC")
        else:
            st.write("  ℹ️ GSC non configuré")

        # ── Claude ──
        st.write("🧠 **Phase 3** — Recommandations Claude...")
        tp_summary = "\n".join([
            f"- {p['url'][:55]} — trafic:{p['traffic']}, KW:{p['top_keyword']} (pos {p['position']})"
            for p in top_pages[:8]
        ]) or "N/A"
        kw_summary = "\n".join([
            f"- {k['keyword']} vol:{k['volume']} pos:{k['position']} mvt:{'↑' + str(k['position_diff']) if k.get('position_diff', 0) > 0 else '↓' + str(abs(k.get('position_diff', 0))) if k.get('position_diff', 0) < 0 else '='}"
            for k in keywords_list[:12]
        ]) or "N/A"

        try:
            reco = generate_recommendations(anthropic_key, {
                "client_name": client_display_name, "domain": domain,
                "current_traffic": fmt(current_traffic), "traffic_delta": traffic_delta,
                "current_keywords": fmt(current_keywords), "keywords_delta": kw_delta,
                "domain_rating": domain_rating, "refdomains": fmt(bl_data["live_refdomains"]),
                "top_pages_summary": tp_summary, "keywords_summary": kw_summary,
                "gsc_clicks": fmt(gsc_report.get("current", {}).get("clicks", 0)) if gsc_report["available"] else "N/A",
                "gsc_impressions": fmt(gsc_report.get("current", {}).get("impressions", 0)) if gsc_report["available"] else "N/A",
                "gsc_ctr": f"{gsc_report.get('current', {}).get('ctr', 0) * 100:.1f}%" if gsc_report["available"] else "N/A",
                "gsc_position": f"{gsc_report.get('current', {}).get('position', 0):.1f}" if gsc_report["available"] else "N/A",
            })
            st.write("  ✅ Recommandations générées")
        except Exception as e:
            reco = ""
            st.write(f"  ⚠️ Erreur Claude : {e}")

        # ── Assemblage PPTX ──
        st.write("🎨 **Phase 4** — Génération du deck Uplix...")

        t_color = UPLIX_GREEN if t_dir == "up" else UPLIX_RED if t_dir == "down" else UPLIX_BLACK
        k_color = UPLIX_GREEN if k_dir == "up" else UPLIX_RED if k_dir == "down" else UPLIX_BLACK

        kpis_list = [
            ("Trafic organique", fmt(current_traffic), traffic_delta, t_color),
            ("Mots-clés", fmt(current_keywords), kw_delta, k_color),
            ("Domain Rating", str(domain_rating or "N/A"), "", UPLIX_BLUE),
            ("Domaines référents", fmt(bl_data["live_refdomains"]), "", UPLIX_TEAL),
        ]

        full_report = {
            "client_name": client_display_name, "domain": domain,
            "period": period_label, "logo_bytes": logo_bytes,
            "kpis_list": kpis_list,
            "traffic_history": traffic_history,
            "top_pages": top_pages, "keywords": keywords_list,
            "backlinks": bl_data, "gsc": gsc_report,
            "recommendations": reco,
        }

        pptx_buf = generate_report_pptx(full_report)
        st.write(f"  ✅ Deck généré — {pptx_buf.getbuffer().nbytes / 1024:.0f} KB")

        status.update(label="Bilan généré !", state="complete", expanded=False)

    # ── Preview ──
    st.divider()
    st.markdown('<div class="section-title">📋 Aperçu du bilan</div>', unsafe_allow_html=True)

    kpi_cols = st.columns(4)
    for i, (label, value, delta, _) in enumerate(kpis_list):
        with kpi_cols[i]:
            st.metric(label, value, delta=delta if delta and delta != "N/A" else None)

    if traffic_history:
        st.subheader("📈 Trafic organique (12 mois)")
        df_h = pd.DataFrame(traffic_history)
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=df_h["date"], y=df_h["traffic"], mode="lines+markers",
                                 line=dict(color=CHART_BLUE, width=3), fill="tozeroy",
                                 fillcolor="rgba(66,133,244,0.1)"))
        fig.update_layout(height=300, margin=dict(t=20, b=40, l=60, r=20))
        st.plotly_chart(fig, use_container_width=True)

    if top_pages:
        st.subheader("🏆 Top pages")
        df_tp = pd.DataFrame(top_pages[:10])[["url", "traffic", "top_keyword", "position"]]
        df_tp.columns = ["URL", "Trafic", "Top mot-clé", "Position"]
        st.dataframe(df_tp, use_container_width=True, hide_index=True)

    if reco:
        st.subheader("💡 Recommandations")
        st.markdown(reco)

    st.divider()

    filename = f"bilan_seo_{domain.replace('.', '_')}_{datetime.now().strftime('%Y%m')}.pptx"
    st.download_button(
        "📥 Télécharger le bilan .pptx",
        data=pptx_buf,
        file_name=filename,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        type="primary",
        use_container_width=True,
    )
    st.caption(f"Bilan généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')} | Importable dans Google Slides")

st.caption("📊 Rapport Client Automatisé | Ma Toolbox SEO")
